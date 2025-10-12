import os
from typing import Tuple
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from langchain_openai import AzureOpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import TokenTextSplitter
import tiktoken
import chardet

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".html"}

def detect_file_encoding(file_path: str) -> str:
    """Detect the encoding of a text file."""
    with open(file_path, "rb") as f:
        raw = f.read(5000)
    result = chardet.detect(raw)
    return result.get("encoding", "utf-8")

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file."""
    text = []
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
    except Exception as e:
        raise RuntimeError(f"Failed to process PDF: {e}")
    return "\n".join(text).strip()

def extract_pptx_text(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file."""
    try:
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs]).strip()
    except Exception as e:
        raise RuntimeError(f"Failed to process DOCX: {e}")

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from TXT file."""
    encoding = detect_file_encoding(file_path)
    with open(file_path, "r", encoding=encoding, errors="ignore") as f:
        return f.read().strip()

def extract_text_from_html(file_path: str) -> str:
    """Extract text from HTML file."""
    from bs4 import BeautifulSoup
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f, "html.parser")
    return soup.get_text(separator="\n").strip()

def extract_text(file_path: str) -> str:
    """
    Unified function to extract text from supported file types.
    
    Args:
        file_path: Path to the file to extract text from
        
    Returns:
        Extracted text content
        
    Raises:
        ValueError: If file format is not supported
        RuntimeError: If file processing fails
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file format: {ext}. Supported formats: {', '.join(SUPPORTED_EXTENSIONS)}")

    try:
        if ext == ".pdf":
            return extract_text_from_pdf(file_path)
        elif ext == ".docx":
            return extract_text_from_docx(file_path)
        elif ext == ".txt":
            return extract_text_from_txt(file_path)
        elif ext == ".html":
            return extract_text_from_html(file_path)
        elif ext == ".pptx":
            return extract_pptx_text(file_path)
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from {file_path}: {e}")

def get_file_info(file_path: str) -> dict:
    """
    Get information about a file.
    
    Args:
        file_path: Path to the file
        
    Returns:
        Dictionary containing file information
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    stat = os.stat(file_path)
    ext = os.path.splitext(file_path)[1].lower()
    
    return {
        "filename": os.path.basename(file_path),
        "extension": ext,
        "size_bytes": stat.st_size,
        "size_mb": round(stat.st_size / (1024 * 1024), 2),
        "is_supported": ext in SUPPORTED_EXTENSIONS,
        "modified_time": stat.st_mtime
    }

def validate_file_for_processing(file_path: str, max_size_mb: int = 10) -> Tuple[bool, str]:
    """
    Validate if a file can be processed.
    
    Args:
        file_path: Path to the file
        max_size_mb: Maximum allowed file size in MB
        
    Returns:
        Tuple of (is_valid, error_message)
    """
    try:
        file_info = get_file_info(file_path)
        
        if not file_info["is_supported"]:
            return False, f"Unsupported file format: {file_info['extension']}"
        
        if file_info["size_mb"] > max_size_mb:
            return False, f"File too large: {file_info['size_mb']}MB (max: {max_size_mb}MB)"
        
        return True, "File is valid for processing"
        
    except Exception as e:
        return False, str(e)
    

# == Handling Large Resources == 
embeddings = AzureOpenAIEmbeddings(
    api_key= os.getenv("AZURE_OPENAI_KEY"),
    azure_endpoint= os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_version= os.getenv("AZURE_OPENAI_API_VERSION"),
    azure_deployment= os.getenv("AZURE_TEXT_EMBEDDING_DEPLOYMENT")
)

def get_relevant_resource_text(text, question, student_answer):
    splitter = TokenTextSplitter(chunk_size=400, chunk_overlap=50, encoding_name="cl100k_base")
    chunks = splitter.split_text(text)
    vectorDB= FAISS.from_texts(chunks, embedding=embeddings)

    query = f"{question}\n{student_answer}"
    retrieved_text = vectorDB.similarity_search(query, k=3)
    resource_text = " ".join([txt.page_content for txt in retrieved_text])
    return resource_text    

def count_tokens(text, model_name="gpt-5-mini"):
    enc = tiktoken.encoding_for_model(model_name)
    count = len(enc.encode(text))
    return count